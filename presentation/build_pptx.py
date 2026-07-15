from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

A = os.path.join(os.path.dirname(__file__), 'assets') + os.sep
BG = RGBColor(0x0f, 0x17, 0x2a)
INK = RGBColor(0xe8, 0xee, 0xf7)
ACC = RGBColor(0x4f, 0x9d, 0xde)
MUT = RGBColor(0x9f, 0xb0, 0xc8)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def bgfill(s):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG


def tb(s, l, t, w, h):
    b = s.shapes.add_textbox(l, t, w, h)
    b.text_frame.word_wrap = True
    return b, b.text_frame


def setpara(p, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    p.text = text
    p.alignment = align
    r = p.runs[0]
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = 'Segoe UI'


def title(s, text, size=34, top=0.5, color=ACC):
    _, tf = tb(s, Inches(0.7), Inches(top), Inches(12), Inches(1.3))
    setpara(tf.paragraphs[0], text, size, color, bold=True)


def bullets(s, items, left=0.9, top=1.9, width=6.0, size=18):
    _, tf = tb(s, Inches(left), Inches(top), Inches(width), Inches(4.8))
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        setpara(p, '-  ' + it, size, INK)
        p.space_after = Pt(10)


def pic(s, name, left, top, width):
    path = A + name
    if os.path.exists(path):
        s.shapes.add_picture(path, left, top, width=width)


# 1 Title
s = prs.slides.add_slide(blank); bgfill(s)
_, tf = tb(s, Inches(1), Inches(2.2), Inches(11.3), Inches(2))
setpara(tf.paragraphs[0], 'Injection Molding Defect Prediction', 40, INK, True, PP_ALIGN.CENTER)
_, tf = tb(s, Inches(1), Inches(4.0), Inches(11.3), Inches(1))
setpara(tf.paragraphs[0], "Predicting good vs scrap parts from a machine's own settings", 20, MUT, False, PP_ALIGN.CENTER)
_, tf = tb(s, Inches(1), Inches(5.3), Inches(11.3), Inches(0.8))
setpara(tf.paragraphs[0], 'James Mahinda', 18, ACC, False, PP_ALIGN.CENTER)

# 2 Problem
s = prs.slides.add_slide(blank); bgfill(s); title(s, 'The problem, and why it matters')
bullets(s, ['Factories catch defective parts late, during inspection',
            'By then material and machine time are already wasted',
            'Every bad cycle is money lost',
            'A drifting machine can scrap a whole batch'])
pic(s, 'machine_photo.jpg', Inches(7.2), Inches(1.9), Inches(5.6))

# 3 How it works
s = prs.slides.add_slide(blank); bgfill(s); title(s, 'How the machine works')
bullets(s, ['Feed plastic pellets in', 'Melt them with a heated screw',
            'Inject into a mold under pressure', 'Cool, then eject the part',
            'One part per cycle, and every cycle its settings are recorded'])
pic(s, 'machine_diagram.png', Inches(7.1), Inches(2.4), Inches(5.8))

# 4 Idea
s = prs.slides.add_slide(blank); bgfill(s); title(s, 'The idea')
bullets(s, ['Use the settings the machine already records',
            'Predict good vs scrap immediately, not after inspection',
            'Point to which settings caused the risk',
            'So an operator fixes the machine before making more scrap'], width=11.5, size=22, top=2.3)

# 5 Data
s = prs.slides.add_slide(blank); bgfill(s); title(s, 'The data')
bullets(s, ['Real production data from a plastics company (iGuzzini)',
            '1,451 molding cycles', '13 process parameters per cycle, plus a quality label',
            'Clean: no missing values, no duplicates'], width=11.5, size=22, top=2.3)

# 6 EDA
s = prs.slides.add_slide(blank); bgfill(s); title(s, 'Do the settings differ between good and scrap?')
bullets(s, ['Yes, several settings clearly differ',
            'Cycle time separates them the most',
            'So there is real signal to learn from'])
pic(s, 'boxplots.png', Inches(7.0), Inches(1.8), Inches(5.9))

# 7 Separability
s = prs.slides.add_slide(blank); bgfill(s); title(s, 'Can a straight line separate them?')
bullets(s, ['Squashed 13 settings into 2D with PCA',
            'Good and scrap overlap, no straight line splits them',
            'Straight line model: scrap F1 = 0.74',
            'Curved model: scrap F1 = 0.85',
            'So the data genuinely needs a flexible model'])
pic(s, 'pca_scatter.png', Inches(7.3), Inches(1.9), Inches(5.5))

# 8 Models
s = prs.slides.add_slide(blank); bgfill(s); title(s, 'The models')
bullets(s, ['Logistic Regression (baseline): scrap F1 = 0.71',
            'SVM (curved): scrap F1 = 0.85',
            'XGBoost: scrap F1 = 0.97, about 98% accuracy',
            'Catches about 71 of 74 scrap, misses only about 3'])
pic(s, 'model_comparison.png', Inches(7.2), Inches(2.0), Inches(5.7))

# 9 Confusion
s = prs.slides.add_slide(blank); bgfill(s); title(s, 'How XGBoost does on unseen parts')
bullets(s, ['291 parts it had never seen',
            'Scrap precision 0.97, recall 0.96',
            'Very few missed defects, very few false alarms'])
pic(s, 'confusion.png', Inches(7.6), Inches(2.0), Inches(4.9))

# 10 SHAP
s = prs.slides.add_slide(blank); bgfill(s); title(s, 'Explaining the prediction (SHAP)')
bullets(s, ['A score alone is not enough on a factory floor',
            'SHAP shows which settings drove each decision',
            'Turns a black box into something an operator can act on'])
pic(s, 'shap_summary.png', Inches(7.2), Inches(1.7), Inches(5.7))

# 11 Anomaly
s = prs.slides.add_slide(blank); bgfill(s); title(s, 'Catching new faults without labels')
bullets(s, ['Isolation Forest learns what a normal cycle looks like',
            'Flags the unusual ones, no labels needed',
            'Catches about a third of scrap on its own',
            'A safety net for factories with no labeled data, or new faults'], width=11.5, size=21, top=2.2)

# 12 Honesty
s = prs.slides.add_slide(blank); bgfill(s); title(s, 'Scope of the results')
bullets(s, ['Scores are measured on a held-out test set the model never saw',
            'Validated on one public dataset: one product, one machine, five production days',
            'Other products or machines would need retraining and fresh validation',
            'The numbers are solid for this data, next is proving them on more machines'], width=11.5, size=20, top=2.2)

# 13 Demo
s = prs.slides.add_slide(blank); bgfill(s)
_, tf = tb(s, Inches(1), Inches(2.6), Inches(11.3), Inches(1.5))
setpara(tf.paragraphs[0], 'Live demo', 44, ACC, True, PP_ALIGN.CENTER)
_, tf = tb(s, Inches(1), Inches(4.1), Inches(11.3), Inches(1.5))
setpara(tf.paragraphs[0], 'Enter a cycle setting, get back the scrap risk and the top settings behind it', 20, MUT, False, PP_ALIGN.CENTER)

# 14 Lessons
s = prs.slides.add_slide(blank); bgfill(s); title(s, 'Lessons learned')
bullets(s, ['Check the simple option first (linear separability) before advanced models',
            'Explaining a model matters as much as its accuracy',
            'Being honest about why a score is high builds trust',
            'The real value is making the output actionable, not just accurate'], width=11.5, size=21, top=2.2)

# 15 Thanks
s = prs.slides.add_slide(blank); bgfill(s); title(s, 'Next steps and thank you')
bullets(s, ['Drift forecasting: predict when the machine is trending toward scrap',
            'Connect to a live machine feed inside a factory system'], width=11.5, size=22, top=2.2)
_, tf = tb(s, Inches(1), Inches(5.2), Inches(11.3), Inches(1))
setpara(tf.paragraphs[0], 'Thank you. Questions?', 24, ACC, False, PP_ALIGN.CENTER)

out = os.path.join(os.path.dirname(__file__), 'slides.pptx')
prs.save(out)
print('Saved', out, 'with', len(prs.slides._sldIdLst), 'slides')
